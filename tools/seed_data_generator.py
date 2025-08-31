#!/usr/bin/env python3
"""
Comprehensive Seed Data Generator for inDoc E2E Testing

This script generates realistic test data for all user roles and document types
to enable comprehensive end-to-end testing and validation.
"""

import asyncio
import sys
import secrets
import logging
import hashlib
import json
import random
from datetime import datetime, timedelta
from pathlib import Path
from typing import Dict, Any, List, Optional
from io import BytesIO
import uuid

# Add the backend directory to the Python path
sys.path.append(str(Path(__file__).parent.parent / "backend"))

from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, delete

from app.db.session import AsyncSessionLocal
from app.models.user import User, UserRole
from app.models.document import Document
from app.models.conversation import Conversation, Message
from app.models.audit import AuditLog
from app.core.security import get_password_hash
from app.core.config import settings

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class SeedDataGenerator:
    """Generate comprehensive seed data for E2E testing"""
    
    def __init__(self, target_total_documents: int = 2000, commit_batch_size: int = 200):
        self.storage_path = Path("../backend/data/storage")
        self.temp_path = Path("../backend/data/temp")
        self.sample_documents = {}
        self.target_total_documents = target_total_documents
        self.commit_batch_size = commit_batch_size
        
        # Ensure directories exist
        self.storage_path.mkdir(parents=True, exist_ok=True)
        self.temp_path.mkdir(parents=True, exist_ok=True)
    
    async def generate_all_seed_data(self, clean_existing: bool = False):
        """Generate all seed data"""
        logger.info("🚀 Starting comprehensive seed data generation...")
        
        async with AsyncSessionLocal() as session:
            if clean_existing:
                await self.clean_existing_data(session)
            
            # Generate users for each role
            users = await self.create_test_users(session)
            
            # Generate diverse document types in bulk
            documents = await self.create_sample_documents(session, users, total=self.target_total_documents)
            
            # Generate conversations and messages
            conversations = await self.create_conversations(session, users, documents)
            
            # Generate audit logs
            await self.create_audit_logs(session, users, documents)
            
            # Generate search scenarios
            await self.create_search_scenarios(session, users, documents)
            
            await session.commit()
            
            logger.info("✅ Seed data generation completed!")
            await self.print_summary(session)
    
    async def clean_existing_data(self, session: AsyncSession):
        """Clean existing test data"""
        logger.info("🧹 Cleaning existing test data...")
        
        # Delete in correct order due to foreign key constraints
        await session.execute(delete(AuditLog))
        await session.execute(delete(Message))
        await session.execute(delete(Conversation))
        await session.execute(delete(Document))
        
        # Keep admin user but delete test users
        await session.execute(
            delete(User).where(User.email.like('%test.indoc.local'))
        )
        
        await session.commit()
        logger.info("✅ Existing test data cleaned")
    
    async def create_test_users(self, session: AsyncSession) -> Dict[str, User]:
        """Create test users for each role with realistic data"""
        logger.info("👥 Creating test users for all roles...")
        
        test_users = {
            # Admin users
            "admin_primary": {
                "email": "admin.primary@test.indoc.local",
                "username": "admin_primary", 
                "full_name": "Alice Administrator",
                "password": secrets.token_urlsafe(16),
                "role": UserRole.ADMIN,
                "is_verified": True
            },
            "admin_secondary": {
                "email": "admin.backup@test.indoc.local",
                "username": "admin_backup",
                "full_name": "Bob Backup Admin", 
                "password": "admin456",
                "role": UserRole.ADMIN,
                "is_verified": True
            },
            
            # Reviewer users
            "reviewer_legal": {
                "email": "legal.reviewer@test.indoc.local",
                "username": "legal_reviewer",
                "full_name": "Carol Legal Reviewer",
                "password": secrets.token_urlsafe(16), 
                "role": UserRole.REVIEWER,
                "is_verified": True
            },
            "reviewer_technical": {
                "email": "tech.reviewer@test.indoc.local",
                "username": "tech_reviewer",
                "full_name": "David Tech Reviewer",
                "password": "review456",
                "role": UserRole.REVIEWER,
                "is_verified": True
            },
            
            # Uploader users
            "uploader_hr": {
                "email": "hr.uploader@test.indoc.local", 
                "username": "hr_uploader",
                "full_name": "Emma HR Uploader",
                "password": secrets.token_urlsafe(16),
                "role": UserRole.UPLOADER,
                "is_verified": True
            },
            "uploader_finance": {
                "email": "finance.uploader@test.indoc.local",
                "username": "finance_uploader", 
                "full_name": "Frank Finance Uploader",
                "password": "upload456",
                "role": UserRole.UPLOADER,
                "is_verified": True
            },
            "uploader_unverified": {
                "email": "new.uploader@test.indoc.local",
                "username": "new_uploader",
                "full_name": "Grace New Uploader",
                "password": "upload789",
                "role": UserRole.UPLOADER,
                "is_verified": False  # Test unverified user workflow
            },
            
            # Viewer users
            "viewer_external": {
                "email": "external.viewer@test.indoc.local",
                "username": "external_viewer",
                "full_name": "Henry External Viewer",
                "password": secrets.token_urlsafe(16),
                "role": UserRole.VIEWER,
                "is_verified": True
            },
            "viewer_intern": {
                "email": "intern.viewer@test.indoc.local", 
                "username": "intern_viewer",
                "full_name": "Ivy Intern Viewer",
                "password": "view456",
                "role": UserRole.VIEWER,
                "is_verified": True
            },
            
            # Compliance users
            "compliance_officer": {
                "email": "officer.compliance@test.indoc.local",
                "username": "compliance_officer",
                "full_name": "Jack Compliance Officer",
                "password": secrets.token_urlsafe(16),
                "role": UserRole.COMPLIANCE,
                "is_verified": True
            },
            "compliance_auditor": {
                "email": "auditor.compliance@test.indoc.local",
                "username": "compliance_auditor", 
                "full_name": "Kate Compliance Auditor",
                "password": "comply456",
                "role": UserRole.COMPLIANCE,
                "is_verified": True
            }
        }
        
        created_users = {}
        
        for user_key, user_data in test_users.items():
            # Check if user already exists
            result = await session.execute(
                select(User).where(User.email == user_data["email"])
            )
            existing_user = result.scalar_one_or_none()
            
            if not existing_user:
                user = User(
                    email=user_data["email"],
                    username=user_data["username"],
                    full_name=user_data["full_name"],
                    hashed_password=get_password_hash(user_data["password"]),
                    role=user_data["role"],
                    is_active=True,
                    is_verified=user_data["is_verified"]
                )
                session.add(user)
                await session.flush()  # Get the ID
                created_users[user_key] = user
                logger.info(f"Created {user_data['role'].value}: {user_data['email']}")
            else:
                created_users[user_key] = existing_user
                logger.info(f"Using existing {user_data['role'].value}: {user_data['email']}")
        
        return created_users
    
    async def create_sample_documents(self, session: AsyncSession, users: Dict[str, User], total: int = 2000) -> List[Document]:
        """Create diverse document types for testing at scale"""
        logger.info("📄 Creating sample documents of various types at scale (target=%s)...", total)
        
        # Document templates with realistic content
        document_templates = [
            # PDF Documents
            {
                "filename": "company_policy_2024.pdf",
                "file_type": "pdf", 
                "title": "Company Policy Manual 2024",
                "description": "Updated company policies and procedures for 2024",
                "tags": ["policy", "hr", "2024", "manual"],
                "content": "COMPANY POLICY MANUAL 2024\n\nEmployee Handbook\n1. Code of Conduct\n2. Remote Work Policy\n3. Data Security Guidelines\n...",
                "access_level": "internal",
                "uploader": "uploader_hr"
            },
            {
                "filename": "financial_report_q4.pdf", 
                "file_type": "pdf",
                "title": "Q4 Financial Report",
                "description": "Quarterly financial performance and analysis",
                "tags": ["finance", "quarterly", "report", "confidential"],
                "content": "QUARTERLY FINANCIAL REPORT Q4 2024\n\nRevenue: $2.5M\nExpenses: $1.8M\nNet Income: $700K\n...",
                "access_level": "confidential",
                "uploader": "uploader_finance"
            },
            {
                "filename": "technical_specification.pdf",
                "file_type": "pdf", 
                "title": "System Architecture Specification",
                "description": "Technical documentation for system architecture",
                "tags": ["technical", "architecture", "specification"],
                "content": "SYSTEM ARCHITECTURE\n\n1. Overview\n2. Components\n3. Data Flow\n4. Security Model\n...",
                "access_level": "internal",
                "uploader": "uploader_hr"
            },
            
            # DOCX Documents
            {
                "filename": "meeting_notes_jan2024.docx",
                "file_type": "docx",
                "title": "January 2024 Board Meeting Notes", 
                "description": "Minutes from the monthly board meeting",
                "tags": ["meeting", "board", "january", "2024"],
                "content": "BOARD MEETING MINUTES\nDate: January 15, 2024\nAttendees: Alice, Bob, Carol\nAgenda Items:\n1. Budget Review\n2. New Initiatives\n...",
                "access_level": "private",
                "uploader": "uploader_hr"
            },
            {
                "filename": "project_proposal.docx",
                "file_type": "docx",
                "title": "AI Integration Project Proposal",
                "description": "Proposal for integrating AI capabilities",
                "tags": ["project", "ai", "proposal", "innovation"],
                "content": "AI INTEGRATION PROJECT PROPOSAL\n\nExecutive Summary\nObjectives\nTimeline\nBudget Requirements\n...",
                "access_level": "internal",
                "uploader": "uploader_finance"
            },
            
            # XLSX Documents
            {
                "filename": "employee_database.xlsx",
                "file_type": "xlsx",
                "title": "Employee Information Database",
                "description": "Comprehensive employee information and records",
                "tags": ["hr", "employees", "database", "confidential"],
                "content": "EMPLOYEE DATABASE\nName\tRole\tDepartment\tSalary\nAlice Admin\tAdmin\tIT\t$85000\nBob Manager\tManager\tHR\t$75000\n...",
                "access_level": "confidential",
                "uploader": "uploader_hr"
            },
            {
                "filename": "budget_analysis_2024.xlsx",
                "file_type": "xlsx",
                "title": "2024 Budget Analysis",
                "description": "Detailed budget breakdown and analysis",
                "tags": ["finance", "budget", "analysis", "2024"],
                "content": "2024 BUDGET ANALYSIS\nDepartment\tAllocated\tSpent\tRemaining\nIT\t$500K\t$300K\t$200K\nHR\t$200K\t$150K\t$50K\n...",
                "access_level": "confidential", 
                "uploader": "uploader_finance"
            },
            
            # PPTX Documents
            {
                "filename": "quarterly_presentation.pptx",
                "file_type": "pptx",
                "title": "Q4 Business Review Presentation",
                "description": "Quarterly business review and future planning",
                "tags": ["presentation", "quarterly", "business", "review"],
                "content": "Q4 BUSINESS REVIEW\nSlide 1: Executive Summary\nSlide 2: Key Metrics\nSlide 3: Achievements\nSlide 4: Challenges\n...",
                "access_level": "internal",
                "uploader": "uploader_finance"
            },
            
            # TXT Documents
            {
                "filename": "api_documentation.txt",
                "file_type": "txt",
                "title": "API Documentation",
                "description": "Technical API reference documentation",
                "tags": ["api", "documentation", "technical", "reference"],
                "content": "API DOCUMENTATION\n\nAuthentication Endpoints:\nPOST /auth/login\nGET /auth/me\n\nDocument Endpoints:\nPOST /files\nGET /files\n...",
                "access_level": "internal",
                "uploader": "uploader_hr"
            },
            {
                "filename": "system_logs.txt",
                "file_type": "txt", 
                "title": "System Error Logs",
                "description": "Recent system error logs for debugging",
                "tags": ["logs", "errors", "system", "debugging"],
                "content": "SYSTEM LOGS\n2024-01-15 10:30:00 ERROR: Database connection timeout\n2024-01-15 10:35:00 INFO: Connection restored\n...",
                "access_level": "private",
                "uploader": "uploader_hr"
            },
            
            # HTML Documents
            {
                "filename": "user_manual.html",
                "file_type": "html",
                "title": "User Manual and Guidelines",
                "description": "Comprehensive user manual for the application",
                "tags": ["manual", "user", "guidelines", "help"],
                "content": "<html><body><h1>User Manual</h1><h2>Getting Started</h2><p>Welcome to inDoc...</p></body></html>",
                "access_level": "public",
                "uploader": "uploader_hr"
            },
            
            # JSON Documents
            {
                "filename": "configuration_settings.json",
                "file_type": "json",
                "title": "System Configuration Settings",
                "description": "Application configuration and settings",
                "tags": ["config", "settings", "system", "technical"],
                "content": '{"database": {"host": "localhost", "port": 5432}, "redis": {"host": "localhost", "port": 6379}}',
                "access_level": "private",
                "uploader": "uploader_hr"
            },
            {
                "filename": "api_responses.json",
                "file_type": "json", 
                "title": "Sample API Response Data",
                "description": "Example API responses for testing",
                "tags": ["api", "responses", "testing", "examples"],
                "content": '{"users": [{"id": 1, "name": "Alice"}, {"id": 2, "name": "Bob"}], "status": "success"}',
                "access_level": "internal",
                "uploader": "uploader_finance"
            }
        ]
        
        created_documents: List[Document] = []

        # Helper: generate real file bytes for some types; fall back to text
        def generate_file_bytes(file_type: str, text_content: str) -> bytes:
            ft = file_type.lower()
            try:
                if ft == "docx":
                    from docx import Document as DocxDocument  # type: ignore
                    docx = DocxDocument()
                    for line in text_content.split("\n"):
                        docx.add_paragraph(line)
                    tmp = BytesIO()
                    docx.save(tmp)
                    return tmp.getvalue()
                if ft == "xlsx":
                    from openpyxl import Workbook  # type: ignore
                    wb = Workbook()
                    ws = wb.active
                    for i, line in enumerate(text_content.split("\n")[:100], start=1):
                        ws.cell(row=i, column=1, value=line)
                    tmp = BytesIO()
                    wb.save(tmp)
                    return tmp.getvalue()
                if ft == "pptx":
                    from pptx import Presentation  # type: ignore
                    prs = Presentation()
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.title.text = "inDoc Presentation"
                    slide.placeholders[1].text = text_content[:200]
                    tmp = BytesIO()
                    prs.save(tmp)
                    return tmp.getvalue()
                if ft in ["png", "jpg", "jpeg"]:
                    from PIL import Image, ImageDraw, ImageFont  # type: ignore
                    img = Image.new("RGB", (640, 360), color=(25, 118, 210))
                    draw = ImageDraw.Draw(img)
                    draw.text((20, 160), "inDoc Sample", fill=(255, 255, 255))
                    tmp = BytesIO()
                    img.save(tmp, format="PNG" if ft == "png" else "JPEG")
                    return tmp.getvalue()
                if ft == "pdf":
                    # Minimal valid PDF with a blank page
                    from pypdf import PdfWriter  # type: ignore
                    writer = PdfWriter()
                    writer.add_blank_page(width=595, height=842)
                    tmp = BytesIO()
                    writer.write(tmp)
                    return tmp.getvalue()
            except Exception:
                pass
            # Fallback: plain text bytes
            return text_content.encode("utf-8")

        # Decide how many documents to create per template
        per_template = max(1, total // len(document_templates))

        batch_counter = 0
        for doc_template in document_templates:
            for i in range(per_template):
                # Vary content a bit to make search meaningful
                varied_content = f"{doc_template['content']}\n\nSeed Index: {i}\nGenerated: {datetime.now().isoformat()}"
                file_bytes = generate_file_bytes(doc_template["file_type"], varied_content)
                file_hash = hashlib.sha256(file_bytes).hexdigest()
                
                # Save file to storage
                storage_path = self.storage_path / f"{file_hash}.{doc_template['file_type']}"
                storage_path.write_bytes(file_bytes)
                
                # Get uploader user (round-robin if missing)
                uploader_key = doc_template.get("uploader")
                if uploader_key not in users:
                    uploader_key = random.choice(list(users.keys()))
                uploader = users[uploader_key]
                
                # Create document record
                document = Document(
                    uuid=uuid.uuid4(),
                    filename=f"{Path(doc_template['filename']).stem}_{i}.{doc_template['file_type']}",
                    file_type=doc_template["file_type"],
                    file_size=len(file_bytes),
                    file_hash=file_hash,
                    storage_path=str(storage_path),
                    status="indexed",
                    virus_scan_status="clean",
                    title=doc_template["title"],
                    description=doc_template["description"],
                    tags=doc_template["tags"],
                    full_text=varied_content,
                    language="en",
                    access_level=doc_template["access_level"],
                    uploaded_by=uploader.id,
                    custom_metadata={
                        "source": "seed_data",
                        "category": doc_template["tags"][0] if doc_template["tags"] else "general",
                        "generated_at": datetime.now().isoformat(),
                        "batch": i // self.commit_batch_size
                    }
                )
                session.add(document)
                created_documents.append(document)
                batch_counter += 1

                if batch_counter % self.commit_batch_size == 0:
                    await session.flush()
                    await session.commit()
                    logger.info("Committed batch of %s documents...", self.commit_batch_size)

        await session.flush()  # Get document IDs for any remaining
        return created_documents
    
    async def create_conversations(self, session: AsyncSession, users: Dict[str, User], documents: List[Document]) -> List[Conversation]:
        """Create realistic conversations and chat scenarios"""
        logger.info("💬 Creating conversations and chat scenarios...")
        
        conversation_scenarios = [
            {
                "title": "Document Search Help",
                "user": "viewer_external",
                "messages": [
                    {"role": "user", "content": "How do I search for financial documents?"},
                    {"role": "assistant", "content": "You can use the search bar to find documents. Try keywords like 'budget', 'financial', or 'report'."},
                    {"role": "user", "content": "Can I filter by document type?"},
                    {"role": "assistant", "content": "Yes! Use the filter options on the left sidebar to filter by file type, date range, or tags."}
                ]
            },
            {
                "title": "Upload Process Inquiry", 
                "user": "uploader_hr",
                "messages": [
                    {"role": "user", "content": "What file formats are supported for upload?"},
                    {"role": "assistant", "content": "We support PDF, DOCX, XLSX, PPTX, TXT, HTML, XML, JSON, EML, and image formats (PNG, JPG, JPEG, TIFF)."},
                    {"role": "user", "content": "Is there a file size limit?"},
                    {"role": "assistant", "content": "Yes, the maximum file size is 100MB per document."}
                ]
            },
            {
                "title": "Compliance Review Process",
                "user": "compliance_officer",
                "messages": [
                    {"role": "user", "content": "How do I review documents for compliance?"},
                    {"role": "assistant", "content": "Navigate to the Audit Trail page to see all document activities. You can filter by user, action, and date range."},
                    {"role": "user", "content": "Can I export audit logs?"},
                    {"role": "assistant", "content": "Yes, use the Export button on the Audit Trail page to download logs in CSV format."}
                ]
            },
            {
                "title": "Technical Support",
                "user": "reviewer_technical", 
                "messages": [
                    {"role": "user", "content": "The search results seem incomplete. What could be wrong?"},
                    {"role": "assistant", "content": "This could be due to indexing delays. Check if Elasticsearch and Weaviate services are running properly."},
                    {"role": "user", "content": "How can I verify the search services?"},
                    {"role": "assistant", "content": "Admin users can check service health on the Settings page under 'System Health'."}
                ]
            }
        ]
        
        created_conversations = []
        
        for scenario in conversation_scenarios:
            user = users[scenario["user"]]
            
            conversation = Conversation(
                title=scenario["title"],
                user_id=user.id,
                tenant_id=uuid.uuid4(),  # For multi-tenancy support
                metadata={
                    "scenario": "e2e_testing",
                    "user_role": user.role.value,
                    "generated_at": datetime.now().isoformat()
                }
            )
            
            session.add(conversation)
            await session.flush()  # Get conversation ID
            
            # Add messages
            for i, msg_data in enumerate(scenario["messages"]):
                message = Message(
                    conversation_id=conversation.id,
                    role=msg_data["role"],
                    content=msg_data["content"],
                    metadata={
                        "message_index": i,
                        "timestamp": (datetime.now() + timedelta(minutes=i)).isoformat()
                    }
                )
                session.add(message)
            
            created_conversations.append(conversation)
            logger.info(f"Created conversation: {scenario['title']} for {user.username}")
        
        return created_conversations
    
    async def create_audit_logs(self, session: AsyncSession, users: Dict[str, User], documents: List[Document]):
        """Create comprehensive audit logs for testing"""
        logger.info("📊 Creating audit logs for compliance testing...")
        
        audit_scenarios = [
            # Document operations
            {"user": "uploader_hr", "action": "upload", "resource_type": "document", "resource_id": None},
            {"user": "reviewer_legal", "action": "view", "resource_type": "document", "resource_id": None},
            {"user": "reviewer_legal", "action": "download", "resource_type": "document", "resource_id": None},
            {"user": "admin_primary", "action": "delete", "resource_type": "document", "resource_id": None},
            
            # User management
            {"user": "admin_primary", "action": "create", "resource_type": "user", "resource_id": None},
            {"user": "admin_primary", "action": "update", "resource_type": "user", "resource_id": None},
            {"user": "admin_secondary", "action": "deactivate", "resource_type": "user", "resource_id": None},
            
            # Search operations
            {"user": "viewer_external", "action": "search", "resource_type": "search", "resource_id": "financial reports"},
            {"user": "viewer_intern", "action": "search", "resource_type": "search", "resource_id": "policy documents"},
            
            # Compliance actions
            {"user": "compliance_officer", "action": "audit_review", "resource_type": "audit", "resource_id": None},
            {"user": "compliance_auditor", "action": "export_logs", "resource_type": "audit", "resource_id": None},
        ]
        
        for i, scenario in enumerate(audit_scenarios):
            user = users[scenario["user"]]
            
            # Use a random document if resource_type is document
            resource_id = scenario["resource_id"]
            if scenario["resource_type"] == "document" and documents:
                resource_id = str(random.choice(documents).id)
            elif scenario["resource_type"] == "user":
                resource_id = str(random.choice(list(users.values())).id)
            
            audit_log = AuditLog(
                user_id=user.id,
                user_email=user.email,
                user_role=user.role.value,
                action=scenario["action"],
                resource_type=scenario["resource_type"],
                resource_id=resource_id,
                ip_address=f"192.168.1.{100 + i}",  # Simulate different IPs
                user_agent="Mozilla/5.0 (Test Browser) inDoc E2E Testing",
                metadata={
                    "test_scenario": True,
                    "scenario_index": i,
                    "generated_at": datetime.now().isoformat()
                }
            )
            
            session.add(audit_log)
        
        logger.info(f"Created {len(audit_scenarios)} audit log entries")
    
    async def create_search_scenarios(self, session: AsyncSession, users: Dict[str, User], documents: List[Document]):
        """Create search test scenarios"""
        logger.info("🔍 Creating search test scenarios...")
        
        # This would typically involve indexing documents in Elasticsearch/Weaviate
        # For now, we'll create metadata that can be used for search testing
        
        search_scenarios = [
            {"query": "financial report", "expected_docs": ["financial_report_q4.pdf"]},
            {"query": "company policy", "expected_docs": ["company_policy_2024.pdf"]}, 
            {"query": "meeting notes", "expected_docs": ["meeting_notes_jan2024.docx"]},
            {"query": "budget analysis", "expected_docs": ["budget_analysis_2024.xlsx"]},
            {"query": "technical specification", "expected_docs": ["technical_specification.pdf"]},
        ]
        
        # Store search scenarios as metadata for testing
        for doc in documents:
            if doc.custom_metadata is None:
                doc.custom_metadata = {}
            
            doc.custom_metadata["search_scenarios"] = [
                scenario for scenario in search_scenarios 
                if doc.filename in scenario.get("expected_docs", [])
            ]
        
        logger.info(f"Created {len(search_scenarios)} search test scenarios")
    
    async def print_summary(self, session: AsyncSession):
        """Print summary of generated data"""
        logger.info("📋 Seed Data Generation Summary:")
        
        # Count users by role
        for role in UserRole:
            result = await session.execute(
                select(User).where(User.role == role)
            )
            count = len(result.scalars().all())
            logger.info(f"  {role.value}: {count} users")
        
        # Count documents by type
        result = await session.execute(select(Document))
        documents = result.scalars().all()
        
        doc_types = {}
        for doc in documents:
            doc_types[doc.file_type] = doc_types.get(doc.file_type, 0) + 1
        
        logger.info("📄 Documents by type:")
        for file_type, count in doc_types.items():
            logger.info(f"  {file_type.upper()}: {count} documents")
        
        # Count conversations
        result = await session.execute(select(Conversation))
        conv_count = len(result.scalars().all())
        logger.info(f"💬 Conversations: {conv_count}")
        
        # Count audit logs
        result = await session.execute(select(AuditLog))
        audit_count = len(result.scalars().all())
        logger.info(f"📊 Audit Logs: {audit_count}")
        
        logger.info("\n🎯 E2E Test Credentials:")
        logger.info("=" * 50)
        
        # Print test user credentials
        test_users = await session.execute(
            select(User).where(User.email.like('%test.indoc.local'))
        )
        
        for user in test_users.scalars().all():
            # Get password from our predefined list
            # Use secure random passwords for all users
            import secrets
            password = secrets.token_urlsafe(16)
            
            logger.info(f"{user.role.value:12} | {user.email:35} | {password}")
        
        logger.info("=" * 50)


async def main():
    """Main function to run seed data generation"""
    # Parse CLI flags
    target_total = 2000
    force_yes = "--yes" in sys.argv or os.getenv("INDOC_YES") == "1"
    flag_clean = "--clean" in sys.argv
    if "--total" in sys.argv:
        try:
            idx = sys.argv.index("--total")
            target_total = int(sys.argv[idx + 1])
        except Exception:
            pass

    generator = SeedDataGenerator(target_total_documents=target_total)
    
    # Ask for confirmation before cleaning existing data unless forced
    if flag_clean:
        logger.warning("⚠️  This will clean existing test data!")
        if not force_yes:
            response = input("Continue? (y/N): ")
            if response.lower() != 'y':
                logger.info("Seed data generation cancelled")
                return
        await generator.generate_all_seed_data(clean_existing=True)
    else:
        await generator.generate_all_seed_data(clean_existing=False)


if __name__ == "__main__":
    asyncio.run(main())
