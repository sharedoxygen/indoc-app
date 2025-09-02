"""
Text extraction service for documents
"""
import logging
import json
import mimetypes
from typing import Optional, Dict, Any, List
from pathlib import Path
import asyncio
from concurrent.futures import ThreadPoolExecutor

# Document processing imports
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    import email
    from email.policy import default
except ImportError:
    email = None

logger = logging.getLogger(__name__)


class TextExtractionService:
    """Service for extracting text from various document formats"""
    
    def __init__(self):
        self.supported_formats = {
            'pdf', 'docx', 'xlsx', 'pptx', 'txt', 
            'html', 'xml', 'json', 'eml', 'csv'
        }
        self.executor = ThreadPoolExecutor(max_workers=4)
    
    async def extract_text(self, file_path: Path) -> Dict[str, Any]:
        """
        Extract text from a document
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dictionary containing extracted text and metadata
        """
        try:
            if not file_path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")
            
            file_extension = file_path.suffix.lower().strip('.')
            
            if file_extension not in self.supported_formats:
                raise ValueError(f"Unsupported file format: {file_extension}")
            
            # Run extraction in thread pool to avoid blocking
            loop = asyncio.get_event_loop()
            result = await loop.run_in_executor(
                self.executor, 
                self._extract_text_sync, 
                file_path, 
                file_extension
            )
            
            logger.info(f"Text extracted from {file_path.name}: {len(result.get('text', ''))} characters")
            return result
            
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            return {
                "text": "",
                "metadata": {
                    "format": file_path.suffix.lower().strip('.'),
                    "error": str(e)
                },
                "success": False,
                "error": str(e)
            }
    
    def _extract_text_sync(self, file_path: Path, file_extension: str) -> Dict[str, Any]:
        """Synchronous text extraction (runs in thread pool)"""
        
        try:
            if file_extension == 'pdf':
                return self._extract_pdf(file_path)
            elif file_extension == 'docx':
                return self._extract_docx(file_path)
            elif file_extension == 'xlsx':
                return self._extract_xlsx(file_path)
            elif file_extension == 'pptx':
                return self._extract_pptx(file_path)
            elif file_extension == 'txt':
                return self._extract_txt(file_path)
            elif file_extension in ['html', 'htm']:
                return self._extract_html(file_path)
            elif file_extension == 'xml':
                return self._extract_xml(file_path)
            elif file_extension == 'json':
                return self._extract_json(file_path)
            elif file_extension == 'eml':
                return self._extract_eml(file_path)
            elif file_extension == 'csv':
                return self._extract_csv(file_path)
            else:
                return self._extract_fallback(file_path)
                
        except Exception as e:
            logger.error(f"Error in sync extraction for {file_path}: {str(e)}")
            raise
    
    def _extract_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from PDF files"""
        if not pypdf:
            raise ImportError("pypdf not available for PDF processing")
        
        text_parts = []
        metadata = {"format": "pdf", "pages": 0}
        
        with open(file_path, 'rb') as file:
            reader = pypdf.PdfReader(file)
            metadata["pages"] = len(reader.pages)
            
            for page_num, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
                except Exception as e:
                    logger.warning(f"Error extracting page {page_num + 1} from {file_path}: {e}")
                    continue
        
        return {
            "text": "\n\n".join(text_parts),
            "metadata": metadata,
            "success": True
        }
    
    def _extract_docx(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from DOCX files"""
        if not DocxDocument:
            raise ImportError("python-docx not available for DOCX processing")
        
        doc = DocxDocument(file_path)
        text_parts = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # Extract table content
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    table_text.append(row_text)
            if table_text:
                text_parts.append("\n[TABLE]\n" + "\n".join(table_text) + "\n[/TABLE]\n")
        
        return {
            "text": "\n".join(text_parts),
            "metadata": {
                "format": "docx",
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables)
            },
            "success": True
        }
    
    def _extract_xlsx(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from XLSX files"""
        if not openpyxl:
            raise ImportError("openpyxl not available for XLSX processing")
        
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text_parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"[SHEET: {sheet_name}]")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip(" |"):
                    text_parts.append(row_text)
            
            text_parts.append("[/SHEET]\n")
        
        return {
            "text": "\n".join(text_parts),
            "metadata": {
                "format": "xlsx",
                "sheets": len(workbook.sheetnames),
                "sheet_names": workbook.sheetnames
            },
            "success": True
        }
    
    def _extract_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from PPTX files"""
        if not Presentation:
            raise ImportError("python-pptx not available for PPTX processing")
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides):
            text_parts.append(f"[SLIDE {slide_num + 1}]")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
            
            text_parts.append("[/SLIDE]\n")
        
        return {
            "text": "\n".join(text_parts),
            "metadata": {
                "format": "pptx",
                "slides": len(prs.slides)
            },
            "success": True
        }
    
    def _extract_txt(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from TXT files"""
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    text = file.read()
                    return {
                        "text": text,
                        "metadata": {
                            "format": "txt",
                            "encoding": encoding,
                            "lines": len(text.splitlines())
                        },
                        "success": True
                    }
            except UnicodeDecodeError:
                continue
        
        # If all encodings fail, read as binary and decode with errors='ignore'
        with open(file_path, 'rb') as file:
            text = file.read().decode('utf-8', errors='ignore')
            return {
                "text": text,
                "metadata": {
                    "format": "txt",
                    "encoding": "utf-8 (with errors ignored)",
                    "lines": len(text.splitlines())
                },
                "success": True
            }
    
    def _extract_html(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from HTML files"""
        if not BeautifulSoup:
            raise ImportError("beautifulsoup4 not available for HTML processing")
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            content = file.read()
        
        soup = BeautifulSoup(content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        text = soup.get_text()
        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = ' '.join(chunk for chunk in chunks if chunk)
        
        return {
            "text": text,
            "metadata": {
                "format": "html",
                "title": soup.title.string if soup.title else None
            },
            "success": True
        }
    
    def _extract_xml(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from XML files"""
        if not BeautifulSoup:
            raise ImportError("beautifulsoup4 not available for XML processing")
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            content = file.read()
        
        soup = BeautifulSoup(content, 'xml')
        text = soup.get_text()
        
        return {
            "text": text,
            "metadata": {
                "format": "xml",
                "root_tag": soup.find().name if soup.find() else None
            },
            "success": True
        }
    
    def _extract_json(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from JSON files"""
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
        
        # Convert JSON to readable text
        text_parts = []
        self._json_to_text(data, text_parts, 0)
        
        return {
            "text": "\n".join(text_parts),
            "metadata": {
                "format": "json",
                "type": type(data).__name__
            },
            "success": True
        }
    
    def _json_to_text(self, obj, text_parts: List[str], depth: int):
        """Recursively convert JSON to readable text"""
        indent = "  " * depth
        
        if isinstance(obj, dict):
            for key, value in obj.items():
                if isinstance(value, (dict, list)):
                    text_parts.append(f"{indent}{key}:")
                    self._json_to_text(value, text_parts, depth + 1)
                else:
                    text_parts.append(f"{indent}{key}: {value}")
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                if isinstance(item, (dict, list)):
                    text_parts.append(f"{indent}[{i}]:")
                    self._json_to_text(item, text_parts, depth + 1)
                else:
                    text_parts.append(f"{indent}[{i}]: {item}")
    
    def _extract_eml(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from EML (email) files"""
        if not email:
            raise ImportError("email module not available for EML processing")
        
        with open(file_path, 'rb') as file:
            msg = email.message_from_binary_file(file, policy=default)
        
        text_parts = []
        
        # Extract headers
        text_parts.append(f"From: {msg.get('From', 'Unknown')}")
        text_parts.append(f"To: {msg.get('To', 'Unknown')}")
        text_parts.append(f"Subject: {msg.get('Subject', 'No Subject')}")
        text_parts.append(f"Date: {msg.get('Date', 'Unknown')}")
        text_parts.append("\n" + "="*50 + "\n")
        
        # Extract body
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    text_parts.append(part.get_content())
        else:
            if msg.get_content_type() == "text/plain":
                text_parts.append(msg.get_content())
        
        return {
            "text": "\n".join(text_parts),
            "metadata": {
                "format": "eml",
                "from": msg.get('From'),
                "subject": msg.get('Subject'),
                "date": msg.get('Date')
            },
            "success": True
        }
    
    def _extract_csv(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from CSV files"""
        import csv
        
        text_parts = []
        row_count = 0
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            # Try to detect dialect
            sample = file.read(1024)
            file.seek(0)
            
            try:
                dialect = csv.Sniffer().sniff(sample)
            except csv.Error:
                dialect = csv.excel
            
            reader = csv.reader(file, dialect)
            
            for row in reader:
                if row:  # Skip empty rows
                    text_parts.append(" | ".join(row))
                    row_count += 1
        
        return {
            "text": "\n".join(text_parts),
            "metadata": {
                "format": "csv",
                "rows": row_count
            },
            "success": True
        }
    
    def _extract_fallback(self, file_path: Path) -> Dict[str, Any]:
        """Fallback extraction for unsupported formats"""
        try:
            # Try to read as text
            return self._extract_txt(file_path)
        except Exception as e:
            return {
                "text": "",
                "metadata": {
                    "format": file_path.suffix.lower().strip('.'),
                    "error": "Unsupported format, could not extract text"
                },
                "success": False,
                "error": f"Unsupported format: {str(e)}"
            }
    
    def extract_text_sync(self, file_path: Path) -> Dict[str, Any]:
        """Synchronous version of extract_text for Celery tasks"""
        import asyncio
        import concurrent.futures
        import threading
        
        # Use thread pool to avoid event loop conflicts
        def run_in_thread():
            # Create new event loop in thread
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            try:
                return loop.run_until_complete(self.extract_text(file_path))
            finally:
                loop.close()
        
        with concurrent.futures.ThreadPoolExecutor() as executor:
            future = executor.submit(run_in_thread)
            return future.result(timeout=300)  # 5 minute timeout